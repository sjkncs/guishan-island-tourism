# -*- coding: utf-8 -*-
"""导出8个月AOI时序数据 + 生成游客画像PPT"""
import json, os, sys
sys.stdout.reconfigure(encoding='utf-8')

try:
    import openpyxl
except ImportError:
    os.system(f'{sys.executable} -m pip install openpyxl')
    import openpyxl

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    os.system(f'{sys.executable} -m pip install python-pptx')
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

# ═══ 数据 ═══
aoi_months = {
    "2025-11": {"uv":4852,"male":0.7578,"female":0.2422,"low":0.1702,"mid":0.4523,"mid_high":0.2794,"high":0.0981,
        "age":{"19-24":0.0498,"25-29":0.0942,"30-34":0.1090,"35-39":0.1544,"40-44":0.1586,"45-49":0.1406,"50-54":0.1567,"55-59":0.0951,"60-64":0.0334,"65+":0.0084}},
    "2025-12": {"uv":4702,"male":0.7792,"female":0.2208,"low":0.1766,"mid":0.4170,"mid_high":0.3051,"high":0.1013,
        "age":{"19-24":0.0461,"25-29":0.0823,"30-34":0.1032,"35-39":0.1477,"40-44":0.1663,"45-49":0.1537,"50-54":0.1602,"55-59":0.0960,"60-64":0.0366,"65+":0.0078}},
    "2026-01": {"uv":4562,"male":0.6727,"female":0.3273,"low":0.1972,"mid":0.4390,"mid_high":0.2699,"high":0.0939,
        "age":{"19-24":0.0906,"25-29":0.1260,"30-34":0.1323,"35-39":0.1464,"40-44":0.1386,"45-49":0.1053,"50-54":0.1333,"55-59":0.0842,"60-64":0.0286,"65+":0.0146}},
    "2026-02": {"uv":4995,"male":0.5834,"female":0.4166,"low":0.1569,"mid":0.4066,"mid_high":0.3132,"high":0.1233,
        "age":{"19-24":0.0500,"25-29":0.0968,"30-34":0.1209,"35-39":0.1743,"40-44":0.1623,"45-49":0.1338,"50-54":0.1307,"55-59":0.0837,"60-64":0.0334,"65+":0.0142}},
    "2026-03": {"uv":4994,"male":0.6526,"female":0.3474,"low":0.1988,"mid":0.3940,"mid_high":0.3091,"high":0.0980,
        "age":{"19-24":0.0460,"25-29":0.0970,"30-34":0.1113,"35-39":0.1591,"40-44":0.1597,"45-49":0.1488,"50-54":0.1556,"55-59":0.0853,"60-64":0.0290,"65+":0.0082}},
    "2026-04": {"uv":5390,"male":0.6432,"female":0.3568,"low":0.1935,"mid":0.4435,"mid_high":0.2679,"high":0.0951,
        "age":{"19-24":0.0479,"25-29":0.0968,"30-34":0.1090,"35-39":0.1633,"40-44":0.1702,"45-49":0.1499,"50-54":0.1495,"55-59":0.0827,"60-64":0.0238,"65+":0.0069}},
    "2026-05": {"uv":5447,"male":0.6340,"female":0.3660,"low":0.1762,"mid":0.4396,"mid_high":0.2801,"high":0.1041,
        "age":{"19-24":0.0504,"25-29":0.0944,"30-34":0.1088,"35-39":0.1636,"40-44":0.1784,"45-49":0.1533,"50-54":0.1411,"55-59":0.0759,"60-64":0.0250,"65+":0.0091}},
    "2026-06": {"uv":4994,"male":0.6601,"female":0.3399,"low":0.1886,"mid":0.4555,"mid_high":0.2666,"high":0.0892,
        "age":{"19-24":0.0428,"25-29":0.1113,"30-34":0.1102,"35-39":0.1542,"40-44":0.1517,"45-49":0.1411,"50-54":0.1444,"55-59":0.0915,"60-64":0.0388,"65+":0.0141}},
}

# ═══ Part 1: Excel ═══
wb = Workbook()
hdr_font = Font(bold=True, color="FFFFFF", size=11, name='Microsoft YaHei')
hdr_fill = PatternFill(start_color="1A2744", end_color="1A2744", fill_type="solid")
body_font = Font(size=10, name='Microsoft YaHei')
thin_border = Border(left=Side(style='thin', color='DDDDDD'), right=Side(style='thin', color='DDDDDD'),
    top=Side(style='thin', color='DDDDDD'), bottom=Side(style='thin', color='DDDDDD'))
title_font = Font(bold=True, size=13, name='Microsoft YaHei')
sub_font = Font(bold=True, size=11, name='Microsoft YaHei', color='1A2744')

months = list(aoi_months.keys())

# Sheet 1: 月度时序
ws1 = wb.active
ws1.title = "AOI月度时序"
ws1.cell(row=1, column=1, value="桂山岛风景区(AOI)月度客流时序 2025.11-2026.06").font = title_font

headers = ["月份", "UV指数", "男性占比", "女性占比", "低消费", "中消费", "中高消费", "高消费"]
for col, h in enumerate(headers, 1):
    c = ws1.cell(row=3, column=col, value=h)
    c.font = hdr_font; c.fill = hdr_fill; c.border = thin_border
for i, m in enumerate(months, 4):
    d = aoi_months[m]
    vals = [m, d["uv"], f'{d["male"]*100:.1f}%', f'{d["female"]*100:.1f}%',
            f'{d["low"]*100:.1f}%', f'{d["mid"]*100:.1f}%', f'{d["mid_high"]*100:.1f}%', f'{d["high"]*100:.1f}%']
    for col, v in enumerate(vals, 1):
        c = ws1.cell(row=i, column=col, value=v)
        c.font = body_font; c.border = thin_border
        if col >= 2: c.alignment = Alignment(horizontal='center')
for col in range(1, 9):
    ws1.column_dimensions[get_column_letter(col)].width = 14

# Sheet 2: 年龄时序
ws2 = wb.create_sheet("AOI年龄时序")
ws2.cell(row=1, column=1, value="桂山岛风景区(AOI)月度年龄分布时序").font = title_font
age_labels = ["19-24","25-29","30-34","35-39","40-44","45-49","50-54","55-59","60-64","65+"]
headers2 = ["年龄段"] + months
for col, h in enumerate(headers2, 1):
    c = ws2.cell(row=3, column=col, value=h)
    c.font = hdr_font; c.fill = hdr_fill; c.border = thin_border
for i, al in enumerate(age_labels, 4):
    ws2.cell(row=i, column=1, value=al).font = body_font
    ws2.cell(row=i, column=1).border = thin_border
    for j, m in enumerate(months, 2):
        v = aoi_months[m]["age"].get(al, 0)
        c = ws2.cell(row=i, column=j, value=f'{v*100:.1f}%')
        c.font = body_font; c.border = thin_border; c.alignment = Alignment(horizontal='center')
for col in range(1, len(headers2)+1):
    ws2.column_dimensions[get_column_letter(col)].width = 12

# Sheet 3: 酒店偏好对比
ws3 = wb.create_sheet("酒店偏好对比")
ws3.cell(row=1, column=1, value="桂山岛AOI vs 香洲区 酒店偏好对比 (2026年6月)").font = title_font
htl_headers = ["酒店类型", "桂山岛AOI", "香洲区", "差异"]
for col, h in enumerate(htl_headers, 1):
    c = ws3.cell(row=3, column=col, value=h)
    c.font = hdr_font; c.fill = hdr_fill; c.border = thin_border

aoi_htl = {"四星级宾馆":0.273263,"旅馆招待所":0.230958,"五星级宾馆":0.226737,"三星级宾馆":0.145958,"经济型连锁酒店":0.117814,"青年旅舍":0.004641,"奢华酒店":0.000629}
dist_htl = {"五星级宾馆":0.326792,"旅馆招待所":0.220127,"四星级宾馆":0.213791,"三星级宾馆":0.136201,"经济型连锁酒店":0.099330,"青年旅舍":0.002651,"奢华酒店":0.001109}
all_htl = sorted(set(list(aoi_htl.keys()) + list(dist_htl.keys())), key=lambda x: -(aoi_htl.get(x,0)))

for i, ht in enumerate(all_htl, 4):
    a = aoi_htl.get(ht, 0); d = dist_htl.get(ht, 0); diff = a - d
    for col, v in enumerate([ht, f'{a*100:.2f}%', f'{d*100:.2f}%', f'{diff*100:+.2f}pp'], 1):
        c = ws3.cell(row=i, column=col, value=v)
        c.font = body_font; c.border = thin_border
        if col >= 2: c.alignment = Alignment(horizontal='center')
for col in range(1, 5):
    ws3.column_dimensions[get_column_letter(col)].width = 18

# Sheet 4: 详细画像
ws4 = wb.create_sheet("AOI详细画像(6月)")
ws4.cell(row=1, column=1, value="桂山岛风景区(AOI)完整画像 2026年6月").font = title_font

sections = [
    ("餐饮消费水平(元)", {"0-49":0.544222,"50-99":0.411096,"100-199":0.039936,"200-299":0.002840,"300-399":0.001157,"400+":0.000748}),
    ("酒店价格偏好(元)", {"0-99":0.487947,"100-299":0.410494,"300-499":0.034231,"500-999":0.038909,"1000+":0.028419}),
    ("长途交通方式", {"火车":0.583172,"飞机":0.302394,"汽车":0.114434}),
    ("日常出行偏好", {"驾车":0.348162,"地铁":0.196021,"火车":0.146529,"打车":0.099215,"飞机":0.064696,"公交":0.051296,"步行":0.050702,"骑行":0.036696}),
]
row = 3
for sec_name, sec_data in sections:
    ws4.cell(row=row, column=1, value=sec_name).font = sub_font
    row += 1
    for k, v in sorted(sec_data.items(), key=lambda x: -x[1]):
        ws4.cell(row=row, column=1, value=k).font = body_font
        ws4.cell(row=row, column=1).border = thin_border
        ws4.cell(row=row, column=2, value=f'{v*100:.2f}%').font = body_font
        ws4.cell(row=row, column=2).border = thin_border
        row += 1
    row += 1
for col in range(1, 3):
    ws4.column_dimensions[get_column_letter(col)].width = 18

OUT = r'C:\Users\Administrator\.qoderwork\workspace\mritsdpuqphp34kx\outputs'
xlsx_path = os.path.join(OUT, '桂山岛_云图时序分析.xlsx')
wb.save(xlsx_path)
import shutil
shutil.copy2(xlsx_path, r'E:\珠海桂山岛案例\数据采集\processed\桂山岛_云图时序分析.xlsx')
print(f"Excel: {xlsx_path}")

# ═══ Part 2: PPT ═══
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
NAVY = RGBColor(0x1A, 0x27, 0x44)
BLUE = RGBColor(0x2E, 0x75, 0xB6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x99, 0x99, 0x99)
DARK = RGBColor(0x33, 0x33, 0x33)

def add_slide(layout_idx=6):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])

def add_text(slide, left, top, width, height, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tf = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height)).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Microsoft YaHei'
    p.alignment = align
    return tf

def add_bullet(slide, left, top, width, height, items, size=14, color=DARK):
    tf = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height)).text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(6)
    return tf

# Slide 1: Title
s = add_slide()
bg = s.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = NAVY
add_text(s, 1, 1.5, 11, 1.5, "桂山岛游客画像分析", size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 1, 3.2, 11, 0.8, "基于高德云图AOI级客流数据 | 2025.11 - 2026.06", size=20, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 1, 4.5, 11, 0.6, "珠海桂山岛风景区(4A) / 香洲区桂山镇 / 城乡规划研究组", size=14, color=GRAY, align=PP_ALIGN.CENTER)

# Slide 2: 数据概览
s = add_slide()
add_text(s, 0.5, 0.3, 12, 0.8, "数据概览", size=32, bold=True, color=NAVY)
metrics = [
    "月均日客流(UV): 约5000 (AOI级, 桂山岛风景区范围内)",
    "月活跃人口(UV): 7462 (月内曾到访桂山岛的人群)",
    "居住人口(UV): 1605 (岛上常住居民)",
    "数据源: 高德云图AOI接口 (B0FFFQ7RVA)",
    "时间范围: 2025年11月 - 2026年6月 (共8个月)",
    "说明: 高德API最早可查2025年11月, 7-10月暑期数据不可获取",
    "UV为人口指数, 非实际人数, 需结合船票数据校准",
]
add_bullet(s, 0.8, 1.3, 11, 5.5, metrics, size=16)

# Slide 3: 客流趋势
s = add_slide()
add_text(s, 0.5, 0.3, 12, 0.8, "客流月度趋势 (2025.11 - 2026.06)", size=32, bold=True, color=NAVY)
trend_items = [
    "UV趋势: 4562(1月低谷) -> 5447(5月高峰), 整体稳定在4700-5400区间",
    "冬季(11-12月): 男性占比高达76-78%, 可能为海钓/户外爱好者",
    "春节(2月): 女性占比跃升至41.7%(全年最高), 家庭出游特征明显",
    "春季(3-5月): 客流回升至5000+, 女性占比稳定在35%左右",
    "6月回落至4994, 进入暑期前过渡期",
]
add_bullet(s, 0.8, 1.3, 11, 5, trend_items, size=16)

# Slide 4: 年龄
s = add_slide()
add_text(s, 0.5, 0.3, 12, 0.8, "年龄结构", size=32, bold=True, color=NAVY)
age_items = [
    "核心客群: 35-54岁占比超60%, 中年家庭客群为绝对主力",
    "其中35-39岁(15-17%)和40-44岁(14-18%)为两个峰值段",
    "50-54岁占14-16%, 中老年客群比例显著高于城市平均",
    "19-24岁仅4-5%(1月异常升至9.1%, 可能元旦假期效应)",
    "25-34岁合计约20%, 青年客群占比偏低",
    "启示: 青年市场(露营/水上运动/社交媒体打卡)存在增长空间",
]
add_bullet(s, 0.8, 1.3, 11, 5.5, age_items, size=16)

# Slide 5: 消费
s = add_slide()
add_text(s, 0.5, 0.3, 12, 0.8, "消费特征", size=32, bold=True, color=NAVY)
cons_items = [
    "消费水平: 中等(45%)+中高(27%)合计超70%, 高消费仅9-12%",
    "餐饮消费: 50元以下占54%, 50-99元占41% -> 岛上餐饮定价应控制在100元以内",
    "酒店消费: 0-99元(49%)+100-299元(41%) -> 民宿/经济型占90%, 高端仅3%",
    "消费季节性: 2月(春节)高消费占比12.3%为全年最高",
    "岛民消费: 低+中占87%, 远低于游客水平, 反映岛上收入结构",
]
add_bullet(s, 0.8, 1.3, 11, 5, cons_items, size=16)

# Slide 6: 酒店对比
s = add_slide()
add_text(s, 0.5, 0.3, 12, 0.8, "酒店偏好对比: 桂山岛 vs 香洲区", size=32, bold=True, color=NAVY)
htl_items = [
    "桂山岛AOI: 四星27% > 旅馆招待所23% > 五星23% > 三星15% > 经济连锁12%",
    "香洲区整体: 五星33% > 旅馆招待所22% > 四星21% > 三星14% > 经济连锁10%",
    "差异: 桂山岛旅馆招待所(民宿)占比高于区级, 五星级低于区级10个百分点",
    "解读: 岛上住宿以特色民宿+品牌酒店为主, 高端度假需求有限",
    "建议: 发展精品民宿而非高端酒店, 四星/民宿是主力供给",
]
add_bullet(s, 0.8, 1.3, 11, 5, htl_items, size=16)

# Slide 7: 交通
s = add_slide()
add_text(s, 0.5, 0.3, 12, 0.8, "交通方式", size=32, bold=True, color=NAVY)
trans_items = [
    "长途交通: 火车58%, 飞机30%, 汽车11%",
    "飞机占比30%高于香洲区整体24% -> 桂山岛吸引较多远程(省外/境外)游客",
    "日常出行: 驾车35%, 地铁20%, 火车15%, 打车10%",
    "地铁占比20%说明大量游客来自有地铁的城市(广深珠等)",
    "岛上无机动车, 步行+骑行是主要出行方式",
    "建议: 优化码头接驳, 增加船班频次, 完善岛上步行系统",
]
add_bullet(s, 0.8, 1.3, 11, 5.5, trans_items, size=16)

# Slide 8: 岛民
s = add_slide()
add_text(s, 0.5, 0.3, 12, 0.8, "岛民画像 vs 游客", size=32, bold=True, color=NAVY)
island_items = [
    "居住人口UV 1605, 活跃人口UV 7462 (游客约为居民4.6倍)",
    "岛民年龄: 40-54岁占47%(近半数), 比游客更老龄化",
    "岛民职业: 商业服务46%(民宿/餐饮/旅游), 公司职员34%",
    "岛民消费: 低+中占87%, 中高+高仅12.5%(游客为36%)",
    "岛上经济高度依赖旅游服务业, 产业结构单一",
    "建议: 发展渔旅融合/文创产品, 增加非住宿类收入来源",
]
add_bullet(s, 0.8, 1.3, 11, 5.5, island_items, size=16)

# Slide 9: 建议
s = add_slide()
add_text(s, 0.5, 0.3, 12, 0.8, "规划建议", size=32, bold=True, color=NAVY)
rec_items = [
    "目标客群: 35-50岁家庭客群(中等消费力), 兼顾青年市场拓展",
    "住宿: 以民宿/精品酒店为主(占需求90%), 控制高端酒店供给",
    "餐饮: 人均控制在100元以内, 发展特色海鲜/渔家体验",
    "交通: 优化珠海-桂山岛船班(火车+飞机为主要到达方式)",
    "青年市场: 海岛露营/水上运动/社交媒体打卡, 提升19-29岁占比",
    "季节性: 冬季开发海钓/观星活动(男性78%), 春季拓展亲子产品",
    "岛民: 发展渔旅融合/文创, 降低对住宿业的过度依赖",
]
add_bullet(s, 0.8, 1.3, 11, 5.5, rec_items, size=16)

# Slide 10: 局限
s = add_slide()
add_text(s, 0.5, 0.3, 12, 0.8, "数据局限性", size=32, bold=True, color=NAVY)
limit_items = [
    "时间范围: AOI数据最早2025年11月, 缺失7-10月暑期旺季数据",
    "行政区接口同样最早2025年11月, 无法回溯更长时间",
    "UV为高德人口指数, 非实际到岛人数, 需结合船票销售数据校准",
    "消费数据基于高德用户行为画像, 不覆盖非高德用户群体",
    "AOI范围为桂山岛风景区, 不含周边岛屿(东澳岛/外伶仃岛等)",
    "建议: 实地调研期间采集船票数据+问卷调查, 与云图数据交叉验证",
]
add_bullet(s, 0.8, 1.3, 11, 5.5, limit_items, size=16)

ppt_path = os.path.join(OUT, '桂山岛游客画像.pptx')
prs.save(ppt_path)
shutil.copy2(ppt_path, r'E:\珠海桂山岛案例\数据采集\processed\桂山岛游客画像.pptx')
print(f"PPT: {ppt_path}")
print("10 slides generated")
